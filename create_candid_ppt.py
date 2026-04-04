#!/usr/bin/env python3
"""
Generate PowerPoint presentation from Candid Studios Media Asset Management Plan
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme
    DARK_BLUE = RGBColor(15, 23, 42)
    ACCENT_BLUE = RGBColor(59, 130, 246)
    LIGHT_GRAY = RGBColor(241, 245, 249)
    TEXT_GRAY = RGBColor(71, 85, 105)

    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = DARK_BLUE
        title_para.alignment = PP_ALIGN.CENTER

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = TEXT_GRAY
            subtitle_para.alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(title, content_items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = DARK_BLUE

        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for item in content_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_GRAY
            p.level = 0
            p.space_after = Pt(12)

        return slide

    def add_two_column_slide(title, left_items, right_items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = DARK_BLUE

        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        for item in left_items:
            p = left_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_GRAY
            p.space_after = Pt(10)

        # Right column
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5.5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True

        for item in right_items:
            p = right_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_GRAY
            p.space_after = Pt(10)

        return slide

    # Slide 1: Title
    add_title_slide(
        "Candid Studios",
        "Multi-Portal Media Asset Management System"
    )

    # Slide 2: Executive Summary
    add_content_slide(
        "Executive Summary",
        [
            "Transform manual file transfer workflows into an intelligent, automated system",
            "Eliminate WeTransfer dependency and 500GB+ manual file management",
            "Enable direct photographer uploads and editor access to 104TB local NAS",
            "Deliver professional client galleries with zero-egress costs",
            "Complete automation through Make.com workflow orchestration"
        ]
    )

    # Slide 3: Current Challenges
    add_content_slide(
        "Current Challenges",
        [
            "❌ Manual WeTransfer uploads creating excessive overhead",
            "❌ Inefficient workflow: Photographer → You → Editor → You → Client",
            "❌ Constant re-downloading and manual file movement",
            "❌ 500GB+ project files requiring high-speed transfers",
            "❌ Multiple user roles without granular access control"
        ]
    )

    # Slide 4: Solution Overview
    add_content_slide(
        "Hybrid Architecture Solution",
        [
            "☁️ Cloudflare R2 - Zero-egress object storage for active projects",
            "💾 104TB Local NAS - Primary archive and editor workspace",
            "🔄 Nextcloud - Collaboration hub with granular access control",
            "⚙️ Make.com - Workflow automation orchestration",
            "🖥️ RaiDrive - Seamless desktop mounting for upload/download",
            "🌐 WordPress + Ultimate Member - Client portals and galleries"
        ]
    )

    # Slide 5: Three-Layer Architecture
    add_content_slide(
        "Three-Layer System Architecture",
        [
            "Layer 1: ORCHESTRATION - Make.com + GoHighLevel",
            "    • Workflow triggers and automation logic",
            "    • Project status tracking and notifications",
            "",
            "Layer 2: COMPUTE/STORAGE - R2, NAS, Nextcloud",
            "    • Cloudflare R2: Active project storage",
            "    • 104TB NAS: Archive and editor workspace",
            "    • Nextcloud: Collaboration platform",
            "",
            "Layer 3: CLIENT ACCESS - Role-based portals",
            "    • Photographers/Videographers: RaiDrive upload",
            "    • Editors: Nextcloud sync or direct NAS access",
            "    • Clients: WordPress gallery portal"
        ]
    )

    # Slide 6: Workflow 1 - Upload
    add_content_slide(
        "Workflow 1: Project Upload",
        [
            "1. Photographer receives project assignment in GoHighLevel",
            "2. Make.com generates temporary R2 credentials (24-48hr expiry)",
            "3. Photographer mounts R2 bucket via RaiDrive as network drive",
            "4. Direct upload to project-specific R2 folder (500GB+ supported)",
            "5. Make.com detects completion and triggers Rclone sync",
            "6. Automated backup: R2 → 104TB NAS with checksum verification",
            "7. GoHighLevel status updated: 'Ready for Editing'"
        ]
    )

    # Slide 7: Workflow 2 - Editor Access
    add_content_slide(
        "Workflow 2: Editor Access",
        [
            "1. Editor assigned in GoHighLevel triggers Make.com workflow",
            "2. Nextcloud Group Folder created with project-specific ACLs",
            "3. Editor receives credentials and Nextcloud Desktop Client",
            "",
            "Remote Editors (India):",
            "    • Nextcloud Desktop Client syncs files locally",
            "    • Work offline, changes auto-sync when online",
            "",
            "Local Editors:",
            "    • Direct SMB/NFS mount to 104TB NAS",
            "    • Immediate access, no sync delay",
            "",
            "4. Editor marks complete, triggers client delivery workflow"
        ]
    )

    # Slide 8: Workflow 3 - Client Delivery
    add_content_slide(
        "Workflow 3: Client Delivery",
        [
            "1. Editor completion triggers archival workflow in Make.com",
            "2. Rclone syncs final files: NAS → R2 delivery bucket",
            "3. Checksum verification ensures data integrity",
            "4. WordPress REST API creates client gallery post",
            "5. Ultimate Member assigns view permissions to specific client",
            "6. Client receives email with gallery link and login credentials",
            "7. Professional gallery experience with zero-egress downloads",
            "8. R2 lifecycle policies: 30 days → Infrequent Access"
        ]
    )

    # Slide 9: Key Technologies - Part 1
    add_two_column_slide(
        "Key Technologies: Storage & Security",
        [
            "Cloudflare R2:",
            "• S3-compatible object storage",
            "• Zero egress fees (vs. $900/mo on AWS)",
            "• 10TB capacity: ~$150/month",
            "• Asia-Pacific location hints",
            "",
            "Cloudflare Worker Security Proxy:",
            "• Path-based access control",
            "• Temporary credential generation",
            "• Project-level authorization",
            "• Workers KV policy store",
            "• SSO integration"
        ],
        [
            "104TB Local NAS:",
            "• Primary archive and editor workspace",
            "• Direct access via SMB/NFS",
            "• RAID 10 configuration",
            "• Existing infrastructure",
            "",
            "Rclone:",
            "• High-performance sync tool",
            "• Checksum verification",
            "• Multipart upload support",
            "• Bandwidth limiting",
            "• Resume on failure"
        ]
    )

    # Slide 10: Key Technologies - Part 2
    add_two_column_slide(
        "Key Technologies: Collaboration & Portals",
        [
            "Nextcloud (Self-Hosted):",
            "• Collaboration hub for editors",
            "• Desktop sync client",
            "• Group Folders with ACLs",
            "• File versioning & comments",
            "• WebDAV mounting",
            "• 8 CPU, 32GB RAM, Redis cache",
            "",
            "RaiDrive:",
            "• Mount R2 as Windows/Mac drive",
            "• Multipart upload support",
            "• Bandwidth management",
            "• $4.84/user/month"
        ],
        [
            "Make.com Automation:",
            "• Workflow orchestration",
            "• GoHighLevel integration",
            "• Email notifications",
            "• Webhook triggers",
            "• Error handling & logging",
            "",
            "WordPress Portals:",
            "• Ultimate Member (7 roles)",
            "• Elementor gallery templates",
            "• WP Media Folder + R2",
            "• Gravity Forms feedback",
            "• Role-based dashboards"
        ]
    )

    # Slide 11: User Roles & Portals
    add_content_slide(
        "Seven Role-Based Portals",
        [
            "1. Photographers - Upload raw files, track project status",
            "2. Videographers - Upload raw footage, equipment checkout",
            "3. Photo Editors - Access projects, upload edited photos",
            "4. Video Editors - Access projects, upload edited videos",
            "5. Vendors - Download approved resources, brand guidelines",
            "6. Affiliates - Referral dashboard, commission tracking",
            "7. Clients - Private galleries, download finals, feedback"
        ]
    )

    # Slide 12: Security Architecture
    add_content_slide(
        "Multi-Layer Security",
        [
            "Cloudflare Access SSO:",
            "    • Google Workspace / Okta integration",
            "    • Per-user authentication",
            "",
            "Path-Based Authorization:",
            "    • Project-specific access control",
            "    • Enforced by Cloudflare Worker",
            "",
            "Temporary Credentials:",
            "    • 24-48 hour expiration",
            "    • Automatic revocation",
            "",
            "Webhook Security:",
            "    • Shared secret tokens",
            "    • IP whitelisting (Make.com static IPs)",
            "    • HTTPS only"
        ]
    )

    # Slide 13: Implementation Roadmap
    add_content_slide(
        "4-Phase Implementation (10 Weeks)",
        [
            "Phase 1 (Weeks 1-2): Foundation & Upload",
            "    • R2 & Rclone setup, GHL integration, RaiDrive deployment",
            "",
            "Phase 2 (Weeks 3-4): Editor Access",
            "    • Nextcloud deployment, Cloudflare Worker security",
            "",
            "Phase 3 (Weeks 5-6): Client Delivery",
            "    • Archival pipeline, WordPress galleries, Ultimate Member",
            "",
            "Phase 4 (Weeks 7-10): Multi-Portal Expansion",
            "    • 7 role-based portals, vendor/affiliate systems"
        ]
    )

    # Slide 14: Cost Analysis
    add_two_column_slide(
        "Monthly Cost Analysis",
        [
            "Cloudflare Services:",
            "• R2 Storage (10TB): $150",
            "• R2 Operations: $4.05",
            "• Workers + KV: $6.50",
            "• Cloudflare Access: $36",
            "Subtotal: $196.55",
            "",
            "Software Licenses:",
            "• RaiDrive Team (10): $48.40",
            "• Make.com Pro: $29",
            "• Nextcloud: $0 (self-hosted)",
            "Subtotal: $77.40"
        ],
        [
            "Infrastructure:",
            "• VPS (8 CPU, 32GB): $40",
            "• Backups: $10",
            "Subtotal: $50",
            "",
            "Total: ~$324/month",
            "",
            "Savings vs. Alternatives:",
            "• WeTransfer: Limited features",
            "• AWS S3 egress: Save $900/mo",
            "• Time savings: ~$1000/mo value",
            "",
            "ROI: Positive in month 1"
        ]
    )

    # Slide 15: Performance Targets
    add_content_slide(
        "Performance & Success Metrics",
        [
            "Technical Performance:",
            "    • 500GB upload: < 4 hours (1 Gbps connection)",
            "    • 100GB download to India: < 3 hours (75 Mbps avg)",
            "    • System uptime: > 99.5%",
            "    • Upload success rate: > 99%",
            "",
            "Efficiency Gains:",
            "    • Manual file transfer time: 4 hrs/day → 30 min/day",
            "    • Project delivery time: 7 days → 3 days average",
            "    • Client satisfaction: 4.2/5 → 4.8/5 target"
        ]
    )

    # Slide 16: Risk Mitigation
    add_content_slide(
        "Disaster Recovery & Data Protection",
        [
            "Three-Tier Backup Strategy:",
            "    1. Primary: 104TB NAS (RAID 10)",
            "    2. Cloud: Cloudflare R2 (30-day active retention)",
            "    3. Offsite: Quarterly external drive backups",
            "",
            "High Availability:",
            "    • NAS: RAID with hot spare drives",
            "    • Nextcloud: Hot standby server option",
            "    • R2: 99.99% SLA from Cloudflare",
            "",
            "Security Monitoring:",
            "    • Daily ClamAV antivirus scans",
            "    • Cloudflare Access audit logs",
            "    • Uptime monitoring for all services"
        ]
    )

    # Slide 17: Key Benefits
    add_content_slide(
        "Transformational Benefits",
        [
            "✅ Zero manual file transfers - Complete automation",
            "✅ Zero egress fees - Unlimited client downloads",
            "✅ Single upload workflow - Photographers upload once",
            "✅ Direct editor access - Work from 104TB NAS (fastest speed)",
            "✅ Professional galleries - Branded client experience",
            "✅ Granular security - Path-based access control",
            "✅ Scalable architecture - Handles 500GB+ projects",
            "✅ Complete audit trail - Track every file movement",
            "✅ 50TB+ capacity - Room for growth",
            "✅ $900/month savings - vs. AWS S3 egress"
        ]
    )

    # Slide 18: Immediate Next Steps
    add_content_slide(
        "Immediate Action Items",
        [
            "This Week:",
            "    ✅ Review and approve plan",
            "    ✅ Provision Cloudflare R2 account",
            "    ✅ Order RaiDrive licenses (start with 5)",
            "    ✅ Schedule kickoff meeting",
            "",
            "Next Week:",
            "    • Begin Phase 1 implementation",
            "    • Install Rclone on NAS server",
            "    • Create test project in GoHighLevel",
            "    • Build first Make.com scenario",
            "",
            "Questions to Resolve:",
            "    • Nextcloud: Self-hosted vs. VPS?",
            "    • SSO provider preference?",
            "    • Initial pilot photographer group?"
        ]
    )

    # Slide 19: Why This Solution Works
    add_content_slide(
        "Critical Success Factors",
        [
            "1. Zero-Egress R2",
            "   Eliminates unpredictable bandwidth costs entirely",
            "",
            "2. Direct NAS Access for Editors",
            "   Fastest possible workflow, no cloud bottleneck",
            "",
            "3. Cloudflare Worker Security",
            "   Granular, path-based access control at edge",
            "",
            "4. Make.com Orchestration",
            "   Ties all systems together seamlessly",
            "",
            "5. Nextcloud Collaboration",
            "   Enterprise-grade file management and sync"
        ]
    )

    # Slide 20: Closing
    add_title_slide(
        "Ready to Transform Your Workflow?",
        "Eliminate the file transfer bottleneck forever"
    )

    # Save presentation
    output_path = '/mnt/m/code/vidismart/Candid-Studios-Media-Asset-Management.pptx'
    prs.save(output_path)
    print(f"✅ Presentation saved: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
