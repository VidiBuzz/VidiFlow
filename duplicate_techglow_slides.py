#!/usr/bin/env python3
"""
Use TechGlow as base, duplicate slides, and replace text with Candid-Cloud content
This preserves TechGlow's actual visual design
"""

from pptx import Presentation
from copy import deepcopy
import shutil

def extract_candid_content(pptx_path):
    """Extract structured content from Candid-Cloud"""
    prs = Presentation(pptx_path)
    slides_content = []

    print(f"Extracting content from: {pptx_path}")
    print(f"Total slides: {len(prs.slides)}\n")

    for i, slide in enumerate(prs.slides):
        slide_data = {'texts': []}

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_data['texts'].append(shape.text.strip())

        slides_content.append(slide_data)
        title = slide_data['texts'][0] if slide_data['texts'] else '(empty)'
        print(f"Slide {i+1}: {title[:60]}")

    return slides_content

def duplicate_slide(prs, slide_index):
    """Duplicate a slide in the presentation"""
    # Get the source slide
    source_slide = prs.slides[slide_index]

    # Get the layout from source slide
    slide_layout = source_slide.slide_layout

    # Create new slide with same layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Copy all shapes from source to new slide
    for shape in source_slide.shapes:
        el = shape.element
        newel = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return new_slide

def replace_text_in_slide(slide, new_texts):
    """Replace text in shapes while preserving formatting"""
    text_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text')]

    # Replace text in order
    for i, shape in enumerate(text_shapes):
        if i < len(new_texts):
            # Preserve the text frame but replace content
            if hasattr(shape, 'text_frame'):
                # Clear existing paragraphs except first
                tf = shape.text_frame
                while len(tf.paragraphs) > 1:
                    tf._element.remove(tf.paragraphs[-1]._element)

                # Update first paragraph
                if tf.paragraphs:
                    tf.paragraphs[0].text = new_texts[i]
            else:
                shape.text = new_texts[i]

def main():
    # Step 1: Extract Candid content
    print("=" * 70)
    print("STEP 1: Extracting content from Candid-Cloud.v1.pptx")
    print("=" * 70)
    candid_content = extract_candid_content('/mnt/m/+Proj/VidiSmart/Candid-Cloud.v1.pptx')

    # Step 2: Copy TechGlow as base
    print("\n" + "=" * 70)
    print("STEP 2: Loading TechGlow.pptx as base")
    print("=" * 70)
    output_path = '/mnt/m/code/vidismart/Candid-with-TechGlow-Design.pptx'
    shutil.copy('/mnt/m/+Proj/VidiSmart/TechGlow.pptx', output_path)

    prs = Presentation(output_path)
    original_slide_count = len(prs.slides)
    print(f"TechGlow base slides: {original_slide_count}")
    print(f"Candid content slides needed: {len(candid_content)}")

    # Step 3: Duplicate slides to match Candid count
    print("\n" + "=" * 70)
    print("STEP 3: Duplicating slides to match count")
    print("=" * 70)

    slides_needed = len(candid_content)
    current_slides = original_slide_count

    # Duplicate the last slide until we have enough
    while current_slides < slides_needed:
        source_index = (current_slides - 1) % original_slide_count
        print(f"Duplicating slide {source_index + 1} -> new slide {current_slides + 1}")
        duplicate_slide(prs, source_index)
        current_slides += 1

    # Step 4: Replace text content
    print("\n" + "=" * 70)
    print("STEP 4: Replacing text with Candid-Cloud content")
    print("=" * 70)

    for i, (slide, content) in enumerate(zip(prs.slides, candid_content)):
        print(f"Updating slide {i+1}...")
        replace_text_in_slide(slide, content['texts'])

    # Save
    prs.save(output_path)

    print("\n" + "=" * 70)
    print("COMPLETE!")
    print("=" * 70)
    print(f"✅ Presentation saved: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Preserves TechGlow visual design with Candid-Cloud content")

if __name__ == "__main__":
    main()
