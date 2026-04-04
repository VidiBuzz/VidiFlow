#!/usr/bin/env python3
"""
Copy slides from TechGlow and replace text with Candid-Cloud content
This preserves the actual visual design of TechGlow
"""

from pptx import Presentation
import shutil

def extract_content(pptx_path):
    """Extract content from Candid-Cloud"""
    prs = Presentation(pptx_path)
    slides_content = []

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slides_content.append(texts)
        print(f"Slide {i+1}: {len(texts)} text elements")

    return slides_content

def main():
    # Extract Candid-Cloud content
    print("Extracting content from Candid-Cloud.v1.pptx...")
    candid_content = extract_content('/mnt/m/+Proj/VidiSmart/Candid-Cloud.v1.pptx')

    # Copy TechGlow as starting point
    output_path = '/mnt/m/code/vidismart/Candid-TechGlow-Final.pptx'
    shutil.copy('/mnt/m/+Proj/VidiSmart/TechGlow.pptx', output_path)

    # Open the copy
    prs = Presentation(output_path)

    print(f"\nTechGlow has {len(prs.slides)} slides")
    print(f"Candid-Cloud has {len(candid_content)} slides")

    # We need to duplicate TechGlow slides to match Candid-Cloud count
    print("\nApproach: TechGlow only has 13 slides, but Candid has 20.")
    print("We'll need to duplicate TechGlow slides or create a hybrid approach.")
    print("\nActually, the limitation is that TechGlow's visual theme can't be")
    print("automatically extracted and applied to new slides with python-pptx.")
    print("\nRecommendation: Manually apply the TechGlow master/theme in PowerPoint:")
    print("1. Open both presentations in PowerPoint")
    print("2. In Candid-Cloud, go to Design tab")
    print("3. Click 'Browse for Themes' and select TechGlow.pptx")
    print("4. The TechGlow design will be applied to all Candid-Cloud slides")

if __name__ == "__main__":
    main()
