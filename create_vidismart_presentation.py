#!/usr/bin/env python3
"""
Create VidiSmart PowerPoint Presentation
Generates a professional deck showcasing the VidiSmart AI Site Planner
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_vidismart_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme (VidiSmart branding)
    DARK_BG = RGBColor(16, 24, 32)  # #101820
    ACCENT_GREEN = RGBColor(93, 235, 90)  # #5deb5a
    YELLOW = RGBColor(251, 191, 36)  # #FBBF24
    GRAY = RGBColor(156, 163, 175)  # text-gray-400
    WHITE = RGBColor(255, 255, 255)

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add dark background
    background = slide1.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "VidiSmart"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(72)
    title_p.font.bold = True
    title_p.font.color.rgb = ACCENT_GREEN
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI Site Planner"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(48)
    subtitle_p.font.color.rgb = WHITE
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.6))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Configure your modern, AI SmartChannel Tech Stack"
    tagline_p = tagline_frame.paragraphs[0]
    tagline_p.font.size = Pt(24)
    tagline_p.font.color.rgb = GRAY
    tagline_p.alignment = PP_ALIGN.CENTER

    # Brand items
    brands_box = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    brands_frame = brands_box.text_frame
    brands_frame.text = "VidiCRM™  •  VidiBlast & VidiMail™  •  VidiBot™"
    brands_p = brands_frame.paragraphs[0]
    brands_p.font.size = Pt(18)
    brands_p.font.color.rgb = YELLOW
    brands_p.alignment = PP_ALIGN.CENTER

    # Slide 2: Overview
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    background2 = slide2.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background2.fill.solid()
    background2.fill.fore_color.rgb = DARK_BG
    background2.line.fill.background()

    # Title
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title2_frame = title2.text_frame
    title2_frame.text = "What is VidiSmart?"
    title2_p = title2_frame.paragraphs[0]
    title2_p.font.size = Pt(44)
    title2_p.font.bold = True
    title2_p.font.color.rgb = ACCENT_GREEN

    # Content
    content2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content2_frame = content2.text_frame
    content2_frame.word_wrap = True

    overview_text = """VidiSmart is an AI-powered tech stack planning application that helps you configure modern, AI-enabled technology stacks for your projects.

Key Features:
• 14 configuration categories across 3 main pillars
• 150+ technology options to choose from
• AI-powered analysis and recommendations
• Real-time 3D visualizations
• Comprehensive tech stack guidance

Perfect for:
• Developers planning new projects
• Product managers scoping technical requirements
• Agencies proposing client solutions
• Startups building their tech foundation"""

    content2_frame.text = overview_text
    for paragraph in content2_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(12)

    # Slide 3: Data Infrastructure
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    background3 = slide3.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background3.fill.solid()
    background3.fill.fore_color.rgb = DARK_BG
    background3.line.fill.background()

    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title3_frame = title3.text_frame
    title3_frame.text = "Pillar 1: Data Infrastructure"
    title3_p = title3_frame.paragraphs[0]
    title3_p.font.size = Pt(44)
    title3_p.font.bold = True
    title3_p.font.color.rgb = ACCENT_GREEN

    content3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content3_frame = content3.text_frame
    content3_frame.word_wrap = True

    infra_text = """Backend & Database Options:
• Supabase (Postgres) - Default recommendation
• MySQL, MongoDB (NoSQL), Airtable
• Self-Hosted Postgres + Node.js
• Laravel (PHP), Firebase (NoSQL)

Content Management (CMS):
• Strapi - Headless CMS (Default)
• Payload CMS, Headless WordPress
• October CMS, Statamic, Winter CMS (Laravel-based)
• Hygraph

Frontend Framework:
• Next.js (React Framework) - Default
• React.js, Nuxt.js (Vue Framework)
• SvelteKit, Angular"""

    content3_frame.text = infra_text
    for paragraph in content3_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(10)

    # Slide 4: AI Transactions & LLM
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    background4 = slide4.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background4.fill.solid()
    background4.fill.fore_color.rgb = DARK_BG
    background4.line.fill.background()

    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title4_frame = title4.text_frame
    title4_frame.text = "Pillar 2: AI Transactions & LLM"
    title4_p = title4_frame.paragraphs[0]
    title4_p.font.size = Pt(44)
    title4_p.font.bold = True
    title4_p.font.color.rgb = ACCENT_GREEN

    content4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content4_frame = content4.text_frame
    content4_frame.word_wrap = True

    ai_text = """Vector Database:
• Pinecone - Default recommendation
• Weaviate, Milvus, Qdrant
• ChromaDB, PGVector
• Elasticsearch Vector Search

Agentified Agent Army:
• n8n - Workflow automation (Default)
• Zapier, Make, LangGraph
• LangChain, Lindy, Relevance AI

CRM & Marketing:
• HubSpot - Default
• GoHighLevel - All-in-one platform
• Salesforce - Enterprise solution"""

    content4_frame.text = ai_text
    for paragraph in content4_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(10)

    # Slide 5: AI Video UI Design
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    background5 = slide5.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background5.fill.solid()
    background5.fill.fore_color.rgb = DARK_BG
    background5.line.fill.background()

    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title5_frame = title5.text_frame
    title5_frame.text = "Pillar 3: AI Video UI Design"
    title5_p = title5_frame.paragraphs[0]
    title5_p.font.size = Pt(44)
    title5_p.font.bold = True
    title5_p.font.color.rgb = ACCENT_GREEN

    content5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content5_frame = content5.text_frame
    content5_frame.word_wrap = True

    video_text = """Project Archetype:
• Video-First Creator Hub - Default
• Knowledge-Based Discussion Forum
• E-learning Community
• Brand-Sponsored Community

AI Video Creation Tools:
• HeyGen - AI avatar videos (Default)
• Tavus, Pictory AI, InVideo
• Google Veo

AI Site Creation:
• Cursor - AI code editor (Default)
• VS Code, Lovable, Windsurf
• Webstudio, GrapesJS, DataButton"""

    content5_frame.text = video_text
    for paragraph in content5_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(10)

    # Slide 6: Multi-Select Categories
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    background6 = slide6.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background6.fill.solid()
    background6.fill.fore_color.rgb = DARK_BG
    background6.line.fill.background()

    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title6_frame = title6.text_frame
    title6_frame.text = "Additional Configuration Options"
    title6_p = title6_frame.paragraphs[0]
    title6_p.font.size = Pt(44)
    title6_p.font.bold = True
    title6_p.font.color.rgb = ACCENT_GREEN

    content6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5.5))
    content6_frame = content6.text_frame
    content6_frame.word_wrap = True

    left_text = """Hosting & Deployment:
• Vercel, Netlify, Railway
• AWS, Google Cloud, Azure
• DigitalOcean, Cloudflare

Primary AI Models:
• Claude (Anthropic)
• GPT-4 (OpenAI)
• Gemini (Google)
• Llama (Meta)

UI/UX Design Tools:
• Figma, Adobe XD
• Sketch, Framer
• Canva, Penpot"""

    content6_frame.text = left_text
    for paragraph in content6_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(8)

    content6b = slide6.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(5.5))
    content6b_frame = content6b.text_frame
    content6b_frame.word_wrap = True

    right_text = """Video Hosting Channel:
• YouTube, Vimeo, Wistia
• Cloudflare Stream
• Twelve Labs AI

Payment Gateways:
• Stripe, PayPal
• Square, Authorize.net
• Razorpay

MCP Servers:
• Filesystem, GitHub, GitLab
• PostgreSQL, SQLite
• Google Drive, Gmail
• Memory, Figma"""

    content6b_frame.text = right_text
    for paragraph in content6b_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(8)

    # Slide 7: Performance Metrics
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    background7 = slide7.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background7.fill.solid()
    background7.fill.fore_color.rgb = DARK_BG
    background7.line.fill.background()

    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title7_frame = title7.text_frame
    title7_frame.text = "Performance Metrics"
    title7_p = title7_frame.paragraphs[0]
    title7_p.font.size = Pt(44)
    title7_p.font.bold = True
    title7_p.font.color.rgb = ACCENT_GREEN

    content7 = slide7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content7_frame = content7.text_frame
    content7_frame.word_wrap = True

    perf_text = """VidiSmart.com delivers exceptional performance:

🚀 Largest Contentful Paint (LCP): 139ms
   • Industry benchmark: < 2500ms = Good
   • VidiSmart: 91% faster than "good" threshold

⚡ Time to First Byte (TTFB): 1ms
   • Nearly instant server response

🎯 Cumulative Layout Shift (CLS): 0.00
   • Perfect score - zero layout shifts
   • Stable, professional user experience

Technical Implementation:
• Canvas-based graphics (Three.js)
• No traditional image loading
• WebGL-accelerated animations
• Optimized JavaScript delivery"""

    content7_frame.text = perf_text
    for paragraph in content7_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(12)

    # Slide 8: Call to Action
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    background8 = slide8.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background8.fill.solid()
    background8.fill.fore_color.rgb = DARK_BG
    background8.line.fill.background()

    # Large centered title
    cta_title = slide8.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    cta_title_frame = cta_title.text_frame
    cta_title_frame.text = "Ready to Build Your Stack?"
    cta_p = cta_title_frame.paragraphs[0]
    cta_p.font.size = Pt(54)
    cta_p.font.bold = True
    cta_p.font.color.rgb = ACCENT_GREEN
    cta_p.alignment = PP_ALIGN.CENTER

    # Visit website
    visit_box = slide8.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    visit_frame = visit_box.text_frame
    visit_frame.text = "Visit: vidismart.com"
    visit_p = visit_frame.paragraphs[0]
    visit_p.font.size = Pt(36)
    visit_p.font.color.rgb = YELLOW
    visit_p.alignment = PP_ALIGN.CENTER

    # Footer text
    footer_box = slide8.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.8))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Helping you to thrive in 2025\n©2025 VidiBuzz - https://VidiBuzz.com"
    for paragraph in footer_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = GRAY
        paragraph.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_path = '/mnt/m/code/vidismart/VidiSmart-AI-Site-Planner.pptx'
    prs.save(output_path)
    print(f"✅ Presentation created successfully: {output_path}")
    print(f"📊 Total slides: 8")
    print(f"🎨 Theme: VidiSmart Brand Colors (Dark mode with green accents)")

if __name__ == "__main__":
    create_vidismart_presentation()
