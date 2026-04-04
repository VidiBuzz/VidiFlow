#!/usr/bin/env python3
"""
Create Google Slides presentation with TechGlow design and Candid-Cloud content
Uses Kumbh Sans font throughout
"""

from pptx import Presentation
import json

# First, extract the Candid-Cloud content
def extract_candid_content():
    prs = Presentation('/mnt/m/+Proj/VidiSmart/Candid-Cloud.v1.pptx')
    slides_content = []

    for i, slide in enumerate(prs.slides):
        slide_data = {
            'number': i + 1,
            'title': '',
            'bullets': []
        }

        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        # First text is usually title
        if texts:
            slide_data['title'] = texts[0]
            slide_data['bullets'] = texts[1:] if len(texts) > 1 else []

        slides_content.append(slide_data)

    return slides_content

# Generate the Google Slides API commands
def generate_google_slides_structure(slides_content):
    """Generate the structure for Google Slides creation"""

    slides_structure = {
        'title': 'Candid Studios - Media Asset Management',
        'theme': {
            'colors': {
                'primary': '#690be5',  # Purple
                'background': '#000000',  # Black
                'text': '#FFFFFF'  # White text
            },
            'fonts': {
                'heading': 'Kumbh Sans',
                'body': 'Kumbh Sans'
            }
        },
        'slides': []
    }

    for slide_data in slides_content:
        slide_obj = {
            'slide_number': slide_data['number'],
            'title': slide_data['title'],
            'content': slide_data['bullets']
        }
        slides_structure['slides'].append(slide_obj)

    return slides_structure

def main():
    print("=" * 70)
    print("CREATING GOOGLE SLIDES WITH TECHGLOW DESIGN")
    print("=" * 70)

    # Extract content
    print("\nExtracting content from Candid-Cloud.v1.pptx...")
    slides_content = extract_candid_content()
    print(f"Extracted {len(slides_content)} slides")

    # Generate structure
    print("\nGenerating Google Slides structure...")
    slides_structure = generate_google_slides_structure(slides_content)

    # Save structure for reference
    output_path = '/mnt/m/code/vidismart/google_slides_structure.json'
    with open(output_path, 'w') as f:
        json.dump(slides_structure, f, indent=2)

    print(f"Structure saved: {output_path}")

    # Print instructions
    print("\n" + "=" * 70)
    print("GOOGLE SLIDES CREATION INSTRUCTIONS")
    print("=" * 70)
    print("""
To create the presentation:

1. The MCP google-workspace tools should be used to:
   - Create a new presentation
   - Set theme with Kumbh Sans font
   - Add 20 slides with the content
   - Apply purple (#690be5) accent color

2. Background images are available in:
   /mnt/m/code/vidismart/techglow_design_extracted/

3. Key settings:
   - Font: Kumbh Sans (all text)
   - Colors: Purple #690be5, Black, White
   - 20 slides total

The structure is ready in google_slides_structure.json
    """)

    # Print first few slides as example
    print("\nFirst 3 slides preview:")
    for slide in slides_structure['slides'][:3]:
        print(f"\nSlide {slide['slide_number']}: {slide['title']}")
        if slide['content']:
            for item in slide['content'][:2]:
                print(f"  • {item[:60]}...")

if __name__ == "__main__":
    main()
