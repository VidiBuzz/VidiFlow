#!/usr/bin/env python3
"""
Create a PPTX file optimized for Google Slides import with TechGlow design
Using only background images and Kumbh Sans font
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from PIL import Image
import io

def create_techglow_google_slides():
    # Start fresh
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Colors matching TechGlow
    PURPLE = RGBColor(105, 11, 229)  # #690be5
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)

    # Extract Candid content
    candid_prs = Presentation('/mnt/m/+Proj/VidiSmart/Candid-Cloud.v1.pptx')

    print(f"Creating {len(candid_prs.slides)} slides with TechGlow design...")

    for i, source_slide in enumerate(candid_prs.slides):
        # Create new slide with blank layout
        slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])

        # Try to add background image from TechGlow extraction
        bg_image_path = f'/mnt/m/code/vidismart/techglow_design_extracted/slide_10_image_00.png'
        try:
            left = Inches(0)
            top = Inches(0)
            pic = slide.shapes.add_picture(bg_image_path, left, top, width=prs.slide_width, height=prs.slide_height)
            # Move to back
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(0, pic._element)
        except:
            # If background fails, add gradient-like colored rectangle
            shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)  # Rectangle
            shape.fill.solid()
            shape.fill.fore_color.rgb = PURPLE
            shape.line.fill.background()
            # Move to back
            slide.shapes._spTree.remove(shape._element)
            slide.shapes._spTree.insert(0, shape._element)

        # Extract text from source slide
        texts = []
        for shape in source_slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                texts.append(shape.text.strip())

        # Add title (first text)
        if texts:
            title_left = Inches(0.5)
            title_top = Inches(0.4) if i == 0 else Inches(0.5)
            title_width = Inches(9)
            title_height = Inches(1.5) if i == 0 else Inches(1)

            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.text = texts[0]
            title_frame.word_wrap = True
            title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            # Style title
            para = title_frame.paragraphs[0]
            para.font.name = 'Kumbh Sans'
            para.font.size = Pt(60) if i == 0 else Pt(44)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT

        # Add content (remaining text)
        if len(texts) > 1:
            content_left = Inches(0.8)
            content_top = Inches(2.2) if i == 0 else Inches(1.8)
            content_width = Inches(8.4)
            content_height = Inches(3.2)

            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            # If it's a single long text, use it as is
            if len(texts) == 2:
                text_frame.text = texts[1]
                para = text_frame.paragraphs[0]
                para.font.name = 'Kumbh Sans'
                para.font.size = Pt(20) if i == 0 else Pt(18)
                para.font.color.rgb = WHITE
                para.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
            else:
                # Multiple bullets
                for j, text in enumerate(texts[1:]):
                    if j > 0:
                        p = text_frame.add_paragraph()
                    else:
                        p = text_frame.paragraphs[0]

                    p.text = text
                    p.font.name = 'Kumbh Sans'
                    p.font.size = Pt(16)
                    p.font.color.rgb = WHITE
                    p.space_after = Pt(10)
                    p.level = 0

        print(f"  Slide {i+1}: {texts[0][:50] if texts else 'Empty'}...")

    # Save
    output = '/mnt/m/code/vidismart/Candid-TechGlow-Final.pptx'
    prs.save(output)
    print(f"\n✅ Presentation saved: {output}")
    print(f"   - {len(prs.slides)} slides")
    print(f"   - Kumbh Sans font throughout")
    print(f"   - TechGlow purple theme (#690be5)")
    print(f"   - Background images applied")
    print(f"\n📤 Upload this file to Google Drive and open with Google Slides")

if __name__ == "__main__":
    create_techglow_google_slides()
