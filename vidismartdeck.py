import collections
import collections.abc
import sys
import os

# =================================================================
# COMPATIBILITY FIX FOR PYTHON 3.10+
# The 'pptx' library expects some things to be in 'collections'
# but they were moved to 'collections.abc'. This fix bridges them.
# =================================================================
try:
    if not hasattr(collections, "Container"):
        collections.Container = collections.abc.Container
    if not hasattr(collections, "Mapping"):
        collections.Mapping = collections.abc.Mapping
    if not hasattr(collections, "MutableMapping"):
        collections.MutableMapping = collections.abc.MutableMapping
    if not hasattr(collections, "Iterable"):
        collections.Iterable = collections.abc.Iterable
    if not hasattr(collections, "Callable"):
        collections.Callable = collections.abc.Callable
except Exception:
    pass

# Try to import the presentation library
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("-" * 60)
    print("ERROR: 'python-pptx' is not installed.")
    print("Please run this command first:")
    print("pip install python-pptx --break-system-packages")
    print("-" * 60)
    sys.exit(1)

def create_deck():
    """Generates the VidiSmart Strategic Presentation."""
    
    prs = Presentation()

    # Standard Slide Layouts
    # 0 = Title Slide
    # 1 = Title and Content
    # 5 = Blank
    
    def add_slide(title_text, points):
        """Helper to create a standard content slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Using blank for custom styling
        
        # Format Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Add Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = 'Arial'
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(0, 0, 0)
        title_p.alignment = PP_ALIGN.LEFT

        # Add Body Text
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        
        for point in points:
            p = body_tf.add_paragraph()
            p.text = point
            p.font.name = 'Arial'
            p.font.size = Pt(20)
            p.space_after = Pt(12)
            p.level = 0
            
        # Global Footer
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
        footer_p = footer_box.text_frame.paragraphs[0]
        footer_p.text = "VidiSmart | James May: 561.510.0021 | https://VidiSmart.com"
        footer_p.font.size = Pt(10)
        footer_p.font.color.rgb = RGBColor(160, 160, 160)
        footer_p.alignment = PP_ALIGN.RIGHT

    # --- SLIDE 1: TITLE SLIDE ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "VidiSmart AI"
    subtitle1.text = "The Architecture of Veracity\nClosing the 'Trust Loop' in Generative AI\n\nJames May\n561.510.0021 | https://VidiSmart.com"
    
    # --- SLIDE 2: THE PROBLEM ---
    add_slide("The Crisis of Confidence", [
        "McKinsey 'Trust Gap': The #1 barrier to enterprise AI adoption.",
        "The 'Stochastic Parrot' Problem: LLMs are probabilistic engines that mimic fluency but lack verifiable truth.",
        "In high-stakes industries (Finance, Legal, Health), 'maybe' is a liability.",
        "Trust is not an emotion; it is an engineered quality of provenance."
    ])

    # --- SLIDE 3: THE TRUST LOOP ---
    add_slide("Closing the Trust Loop", [
        "Traditional GenAI: Query → Generation (Static/Probabilistic).",
        "VidiSmart Trust Loop: Query → Verification → Grounding → Response.",
        "Closing the loop requires deterministic proof for every synthetic claim.",
        "Objective: Moving from 'Generative' to 'Authoritative' AI."
    ])

    # --- SLIDE 4: THE ARCHITECTURE ---
    add_slide("The Architecture of Veracity", [
        "A system designed to eliminate hallucinations through multimodal grounding.",
        "Unifying 5 modalities of truth into a single parallel execution path.",
        "Built on the Vidi.ai SmartStack to ensure real-time reliability.",
        "Scalable, federated, and enterprise-sovereign."
    ])

    # --- SLIDE 5: THE FIVE PILLARS ---
    add_slide("The Five Pillars of Veracity", [
        "1. Vector: Semantic intent and high-level conceptual mapping.",
        "2. Lexical: Keyword precision for exact entity matching.",
        "3. Graph: Relational logic and knowledge-path validation.",
        "4. Geospatial: Physical anchoring to real-world location context.",
        "5. Visual: Perceptual certainty through multi-modal imagery."
    ])

    # --- SLIDE 6: THE SMARTSTACK ---
    add_slide("The Ai-SmartStack", [
        "Powered by Vidi.ai: The convergence of search and reasoning.",
        "Converged execution: One engine for tensors, vectors, and text.",
        "Compute-Over-Data: Processing on the node to eliminate RAG latency.",
        "Scalable Infrastructure: Built for billion-scale production."
    ])

    # --- SLIDE 7: VISUAL VECTOR ---
    add_slide("The Visual Vector Advantage", [
        "Beyond simple vector databases: A true AI Search Platform.",
        "Omni-Hybrid Capability: Simultaneous lexical and semantic search.",
        "On-Node Inference: Moving the model to the data for 10x performance.",
        "Stateful Learning: Adaptive feedback loops that update in real-time."
    ])

    # --- SLIDE 8: ORCHESTRATION ---
    add_slide("Agentic Orchestration", [
        "The 'Brain' is a dynamic reasoning workflow, not a static index.",
        "Agents decompose complex intent into executable sub-tasks.",
        "Federated coordination of Visual Vector, OpenSearch, and Trino.",
        "Final response synthesis based on multiple verified ground truths."
    ])

    # --- SLIDE 9: HYPERLOCAL VISION ---
    add_slide("Hyperlocal AI: VidiCity", [
        "Integrating VidiShop & VidiCommerce for local economic growth.",
        "Hyperlocal Visual Vectoring: Indexing physical assets for instant discovery.",
        "Performance Max Unification: Combining paid search data with visual intent.",
        "Mapping the economic fabric of 138+ high-growth local sectors."
    ])

    # --- SLIDE 10: CONCLUSION ---
    add_slide("The Future of Trusted AI", [
        "VidiSmart moves AI from 'Information Retrieval' to 'Truth Verification'.",
        "An infrastructure built for the next decade of autonomous systems.",
        "Turning 'Stochastic Parrots' into reliable, cited experts.",
        "Because in the enterprise, getting the facts right is the only metric that matters."
    ])

    # --- SLIDE 11: CONTACT ---
    add_slide("Connect with VidiSmart", [
        "James May",
        "561.510.0021",
        "https://VidiSmart.com",
        "Follow VidiSmart for updates on the Architecture of Veracity."
    ])

    # Save
    output_name = "VidiSmart_Presentation.pptx"
    try:
        prs.save(output_name)
        print("\n" + "="*40)
        print(f"SUCCESS! {output_name} has been created.")
        print("Upload this file to Google Drive to open as Google Slides.")
        print("="*40 + "\n")
    except Exception as e:
        print(f"ERROR saving file: {e}")

if __name__ == "__main__":
    create_deck()