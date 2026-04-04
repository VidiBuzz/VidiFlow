#!/usr/bin/env python3
"""
Extract content from TechGlow.pptx and apply Candid-Cloud.v1.pptx styling
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def extract_presentation_content(pptx_path):
    """Extract all text content from a PowerPoint presentation"""
    prs = Presentation(pptx_path)
    slides_content = []

    print(f"Reading: {pptx_path}")
    print(f"Total slides: {len(prs.slides)}\n")

    for i, slide in enumerate(prs.slides):
        slide_data = {
            'slide_number': i + 1,
            'title': '',
            'content': []
        }

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    # Try to identify if it's a title (usually first or largest text)
                    if not slide_data['title'] and len(text) < 150:
                        slide_data['title'] = text
                    else:
                        if text != slide_data['title']:
                            slide_data['content'].append(text)

        slides_content.append(slide_data)
        print(f"Slide {i+1}: {slide_data['title'][:60] if slide_data['title'] else '(No title)'}")

    return slides_content

def extract_template_style(pptx_path):
    """Extract styling information from Candid-Cloud template"""
    prs = Presentation(pptx_path)

    style_info = {
        'slide_width': prs.slide_width,
        'slide_height': prs.slide_height,
        'slide_layouts': len(prs.slide_layouts)
    }

    print(f"\nTemplate info from {pptx_path}:")
    print(f"  Slide dimensions: {prs.slide_width / 914400:.1f}\" x {prs.slide_height / 914400:.1f}\"")
    print(f"  Available layouts: {len(prs.slide_layouts)}")

    return prs, style_info

def apply_style_to_content(template_prs, content_slides, output_path):
    """Create new presentation using template styling with extracted content"""

    # Create new presentation using the template as base
    # This preserves the template's master slides and layouts
    from copy import deepcopy
    import tempfile

    # Copy the template to use as base
    new_prs = Presentation('/mnt/m/+Proj/VidiSmart/Candid-Cloud.v1.pptx')

    # Remove all existing slides from template
    slide_ids = list(new_prs.slides._sldIdLst)
    for slide_id in slide_ids:
        rId = slide_id.rId
        new_prs.part.drop_rel(rId)
        new_prs.slides._sldIdLst.remove(slide_id)

    # Define color scheme (matching Candid-Cloud style)
    DARK_BLUE = RGBColor(15, 23, 42)
    ACCENT_BLUE = RGBColor(59, 130, 246)
    TEXT_GRAY = RGBColor(71, 85, 105)

    def add_slide_with_content(slide_data):
        """Add a slide with the extracted content"""
        slide = new_prs.slides.add_slide(new_prs.slide_layouts[6])  # Blank layout

        # Add title if present
        if slide_data['title']:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = DARK_BLUE

        # Add content
        if slide_data['content']:
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True

            for item in slide_data['content']:
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(14)
                p.font.color.rgb = TEXT_GRAY
                p.space_after = Pt(10)

    print(f"\nCreating {len(content_slides)} slides...")

    for slide_data in content_slides:
        add_slide_with_content(slide_data)

    # Save
    new_prs.save(output_path)
    print(f"\n✅ Presentation saved: {output_path}")
    return output_path

def main():
    techglow_path = '/mnt/m/+Proj/VidiSmart/TechGlow.pptx'
    candid_cloud_path = '/mnt/m/+Proj/VidiSmart/Candid-Cloud.v1.pptx'
    output_path = '/mnt/m/code/vidismart/TechGlow-Candid-Style.pptx'

    # Step 1: Extract content from TechGlow
    print("=" * 60)
    print("STEP 1: Extracting content from TechGlow.pptx")
    print("=" * 60)
    techglow_content = extract_presentation_content(techglow_path)

    # Step 2: Load Candid-Cloud template
    print("\n" + "=" * 60)
    print("STEP 2: Loading Candid-Cloud.v1.pptx template")
    print("=" * 60)
    template_prs, style_info = extract_template_style(candid_cloud_path)

    # Step 3: Apply styling
    print("\n" + "=" * 60)
    print("STEP 3: Applying Candid-Cloud style to TechGlow content")
    print("=" * 60)
    output = apply_style_to_content(template_prs, techglow_content, output_path)

    print("\n" + "=" * 60)
    print("COMPLETE!")
    print("=" * 60)
    print(f"New presentation: {output}")
    print(f"Original slides: {len(techglow_content)}")

if __name__ == "__main__":
    main()
