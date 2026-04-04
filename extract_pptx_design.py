#!/usr/bin/env python3
"""
Extract visual design elements from TechGlow PowerPoint
- Extract background images from each slide
- Analyze colors, gradients, and design patterns
- Save as PNGs for analysis
"""

from pptx import Presentation
from PIL import Image
import io
import os
import json
from collections import Counter

def extract_images_from_slide(slide, slide_num, output_dir):
    """Extract all images from a slide"""
    images = []

    for shape in slide.shapes:
        # Check if shape has an image
        if hasattr(shape, 'image'):
            image_stream = io.BytesIO(shape.image.blob)
            img = Image.open(image_stream)

            img_path = f"{output_dir}/slide_{slide_num:02d}_image_{len(images):02d}.png"
            img.save(img_path)
            images.append({
                'path': img_path,
                'size': img.size,
                'format': img.format
            })
            print(f"  Extracted image: {img_path} ({img.size[0]}x{img.size[1]})")

    return images

def extract_colors_from_presentation(prs):
    """Extract color scheme from presentation"""
    colors = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'fill'):
                fill = shape.fill
                if fill.type == 1:  # Solid fill
                    if hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'rgb'):
                        rgb = fill.fore_color.rgb
                        colors.append(f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")

    return colors

def analyze_slide_background(slide, slide_num):
    """Analyze slide background properties"""
    background_info = {
        'slide_number': slide_num,
        'has_background': False,
        'type': None
    }

    # Check background
    if hasattr(slide, 'background'):
        background_info['has_background'] = True
        bg = slide.background

        if hasattr(bg, 'fill'):
            fill = bg.fill
            if fill.type:
                background_info['type'] = str(fill.type)

    return background_info

def extract_slide_as_image(prs_path, slide_num, output_path):
    """This requires conversion - we'll note it for later"""
    return None

def main():
    techglow_path = '/mnt/m/+Proj/VidiSmart/TechGlow.pptx'
    output_dir = '/mnt/m/code/vidismart/techglow_design_extracted'

    # Create output directory
    os.makedirs(output_dir, exist_ok=True)

    print("=" * 70)
    print("EXTRACTING DESIGN ELEMENTS FROM TECHGLOW")
    print("=" * 70)

    prs = Presentation(techglow_path)

    design_analysis = {
        'slide_count': len(prs.slides),
        'slide_dimensions': {
            'width': prs.slide_width / 914400,  # Convert to inches
            'height': prs.slide_height / 914400
        },
        'slides': [],
        'color_palette': [],
        'fonts': set()
    }

    # Extract from each slide
    print(f"\nAnalyzing {len(prs.slides)} slides...\n")

    all_images = []
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        print(f"Slide {slide_num}:")

        # Extract images
        slide_images = extract_images_from_slide(slide, slide_num, output_dir)
        all_images.extend(slide_images)

        # Analyze background
        bg_info = analyze_slide_background(slide, slide_num)

        # Extract text and fonts
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                slide_texts.append(shape.text.strip())

                # Try to get font info
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if hasattr(run.font, 'name') and run.font.name:
                                design_analysis['fonts'].add(run.font.name)

        design_analysis['slides'].append({
            'number': slide_num,
            'text_count': len(slide_texts),
            'image_count': len(slide_images),
            'background': bg_info
        })

        print(f"  Text elements: {len(slide_texts)}")
        print(f"  Images: {len(slide_images)}")
        print()

    # Extract colors
    print("Extracting color palette...")
    colors = extract_colors_from_presentation(prs)
    color_counts = Counter(colors)
    design_analysis['color_palette'] = [
        {'color': color, 'count': count}
        for color, count in color_counts.most_common(10)
    ]

    # Convert fonts set to list for JSON
    design_analysis['fonts'] = list(design_analysis['fonts'])

    # Save analysis
    analysis_path = f"{output_dir}/design_analysis.json"
    with open(analysis_path, 'w') as f:
        json.dump(design_analysis, f, indent=2)

    print("=" * 70)
    print("EXTRACTION COMPLETE")
    print("=" * 70)
    print(f"\nOutput directory: {output_dir}")
    print(f"Total images extracted: {len(all_images)}")
    print(f"Design analysis saved: {analysis_path}")
    print(f"\nTop colors found:")
    for item in design_analysis['color_palette'][:5]:
        print(f"  {item['color']} (used {item['count']} times)")

    if design_analysis['fonts']:
        print(f"\nFonts used: {', '.join(design_analysis['fonts'])}")

    print("\n" + "=" * 70)
    print("NEXT STEP: Convert slides to PNG for full visual analysis")
    print("=" * 70)
    print("We need to convert the actual slides to images to see the full design.")
    print("This requires LibreOffice or similar conversion tool.")

if __name__ == "__main__":
    main()
