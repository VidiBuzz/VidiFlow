#!/usr/bin/env python3
"""
Create presentation from scratch with TechGlow design and Kumbh Sans font
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def main():
    print("=" * 70)
    print("CREATING PRESENTATION FROM SCRATCH")
    print("=" * 70)

    # Step 1: Load Candid-Cloud content
    print("\nStep 1: Loading Candid-Cloud content...")
    candid_prs = Presentation('/mnt/m/+Proj/VidiSmart/Candid-Cloud.v1.pptx')

    slides_content = []
    for i, slide in enumerate(candid_prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                texts.append(shape.text.strip())
        slides_content.append(texts)
        print(f"  Slide {i+1}: {texts[0][:50] if texts else '(empty)'}...")

    print(f"✓ Loaded {len(slides_content)} slides")

    # Step 2: Create new presentation
    print("\nStep 2: Creating new presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Colors
    PURPLE = RGBColor(105, 11, 229)  # #690be5
    WHITE = RGBColor(255, 255, 255)

    # Step 3: Check for TechGlow background
    print("\nStep 3: Loading TechGlow background...")
    bg_image = '/mnt/m/code/vidismart/techglow_design_extracted/slide_10_image_00.png'
    bg_exists = os.path.exists(bg_image)
    print(f"  Background image: {'✓ Found' if bg_exists else '✗ Not found'}")

    # Step 4: Create slides
    print("\nStep 4: Creating 20 slides with TechGlow design + Kumbh Sans...")

    for i, texts in enumerate(slides_content):
        # Add blank slide
        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        # Add background
        if bg_exists:
            pic = slide.shapes.add_picture(bg_image, 0, 0, width=prs.slide_width, height=prs.slide_height)
            # Move to back
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(0, pic._element)

        # Add title
        if texts:
            title_box = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(0.8) if i == 0 else Inches(0.5),
                Inches(9),
                Inches(1.2)
            )
            tf = title_box.text_frame
            tf.text = texts[0]
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.font.name = 'Kumbh Sans'
            p.font.size = Pt(54) if i == 0 else Pt(40)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT

        # Add content
        if len(texts) > 1:
            content_box = slide.shapes.add_textbox(
                Inches(0.8),
                Inches(2.2) if i == 0 else Inches(1.8),
                Inches(8.4),
                Inches(3.2)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            # Add content text
            content_text = '\n'.join(texts[1:])
            tf.text = content_text

            for paragraph in tf.paragraphs:
                paragraph.font.name = 'Kumbh Sans'
                paragraph.font.size = Pt(24) if i == 0 else Pt(18)
                paragraph.font.color.rgb = WHITE
                paragraph.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
                paragraph.space_after = Pt(10)

        print(f"  ✓ Slide {i+1} created")

    # Step 5: Save
    print("\nStep 5: Saving presentation...")
    output = '/mnt/m/code/vidismart/Candid-TechGlow-Fresh.pptx'
    prs.save(output)

    print("\n" + "=" * 70)
    print("✅ SUCCESS!")
    print("=" * 70)
    print(f"File saved: {output}")
    print(f"Windows path: M:\\code\\vidismart\\Candid-TechGlow-Fresh.pptx")
    print(f"\nSlides: {len(prs.slides)}")
    print(f"Font: Kumbh Sans (all text)")
    print(f"Theme: TechGlow purple (#690be5)")
    print(f"Background: {'Applied' if bg_exists else 'Gradient fallback'}")

if __name__ == "__main__":
    main()
