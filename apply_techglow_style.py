#!/usr/bin/env python3
"""
Extract content from Candid-Cloud.v1.pptx and apply TechGlow.pptx styling
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def extract_presentation_content(pptx_path):
    """Extract all text content from a PowerPoint presentation"""
    prs = Presentation(pptx_path)
    slides_content = []

    print(f"Reading: {pptx_path}")
    print(f"Total slides: {len(prs.slides)}\n")

    for i, slide in enumerate(prs.slides):
        slide_data = {
            'slide_number': i + 1,
            'title': '',
            'content': []
        }

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    # Try to identify if it's a title (usually first or largest text)
                    if not slide_data['title'] and len(text) < 150:
                        slide_data['title'] = text
                    else:
                        if text != slide_data['title']:
                            slide_data['content'].append(text)

        slides_content.append(slide_data)
        print(f"Slide {i+1}: {slide_data['title'][:60] if slide_data['title'] else '(No title)'}")

    return slides_content

def apply_techglow_style(content_slides, output_path):
    """Create new presentation using TechGlow template with Candid-Cloud content"""

    # Load TechGlow template to preserve its master slides and layouts
    new_prs = Presentation('/mnt/m/+Proj/VidiSmart/TechGlow.pptx')

    # Remove all existing slides from template
    slide_ids = list(new_prs.slides._sldIdLst)
    for slide_id in slide_ids:
        rId = slide_id.rId
        new_prs.part.drop_rel(rId)
        new_prs.slides._sldIdLst.remove(slide_id)

    # Define color scheme (matching TechGlow style - vibrant gradients)
    GRADIENT_BLUE = RGBColor(59, 130, 246)
    GRADIENT_PURPLE = RGBColor(139, 92, 246)
    DARK_TEXT = RGBColor(15, 23, 42)
    LIGHT_TEXT = RGBColor(148, 163, 184)

    def add_slide_with_content(slide_data):
        """Add a slide with the extracted content using TechGlow layouts"""
        # TechGlow has only 1 layout, use it
        slide = new_prs.slides.add_slide(new_prs.slide_layouts[0])

        # Add title if present
        if slide_data['title']:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = DARK_TEXT
            title_para.alignment = PP_ALIGN.CENTER

        # Add content
        if slide_data['content']:
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True

            for item in slide_data['content']:
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_TEXT
                p.space_after = Pt(12)

    print(f"\nCreating {len(content_slides)} slides with TechGlow styling...")

    for slide_data in content_slides:
        add_slide_with_content(slide_data)

    # Save
    new_prs.save(output_path)
    print(f"\n✅ Presentation saved: {output_path}")
    return output_path

def main():
    candid_cloud_path = '/mnt/m/+Proj/VidiSmart/Candid-Cloud.v1.pptx'
    output_path = '/mnt/m/code/vidismart/Candid-Cloud-TechGlow-Style.pptx'

    # Step 1: Extract content from Candid-Cloud
    print("=" * 60)
    print("STEP 1: Extracting content from Candid-Cloud.v1.pptx")
    print("=" * 60)
    candid_content = extract_presentation_content(candid_cloud_path)

    # Step 2: Apply TechGlow styling
    print("\n" + "=" * 60)
    print("STEP 2: Applying TechGlow style to Candid-Cloud content")
    print("=" * 60)
    output = apply_techglow_style(candid_content, output_path)

    print("\n" + "=" * 60)
    print("COMPLETE!")
    print("=" * 60)
    print(f"New presentation: {output}")
    print(f"Total slides: {len(candid_content)}")

if __name__ == "__main__":
    main()
