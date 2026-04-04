#!/usr/bin/env python3
"""
Create VidiSmart PowerPoint Presentation - Restructured
Each menu category gets its own slide
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_dark_background(slide, dark_bg, prs):
    """Add dark background to slide"""
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = dark_bg
    background.line.fill.background()

def add_slide_title(slide, title_text, accent_green):
    """Add title to slide"""
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title_frame = title.text_frame
    title_frame.text = title_text
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = accent_green
    return title

def add_option_list(slide, items, white):
    """Add list of options to slide"""
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content_frame = content.text_frame
    content_frame.word_wrap = True

    for item in items:
        p = content_frame.add_paragraph() if content_frame.paragraphs[0].text else content_frame.paragraphs[0]
        p.text = f"• {item['name']}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = white
        p.space_after = Pt(4)

        # Add description
        desc_p = content_frame.add_paragraph()
        desc_p.text = f"  {item['description']}"
        desc_p.font.size = Pt(14)
        desc_p.font.color.rgb = RGBColor(200, 200, 200)
        desc_p.space_after = Pt(12)
        desc_p.level = 1

def create_vidismart_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define colors
    DARK_BG = RGBColor(16, 24, 32)
    ACCENT_GREEN = RGBColor(93, 235, 90)
    YELLOW = RGBColor(251, 191, 36)
    GRAY = RGBColor(156, 163, 175)
    WHITE = RGBColor(255, 255, 255)

    # Technology data from script.js
    tech_data = {
        'Project Archetype': [
            {'name': 'Video-First Creator Hub', 'description': 'A platform focused on video content, community engagement, and monetization for creators.'},
            {'name': 'Knowledge-Based Discussion Forum', 'description': 'A community centered around structured discussions, knowledge sharing, and long-form content.'},
            {'name': 'E-learning Community', 'description': 'An integrated platform for online courses, student interaction, and educational content delivery.'},
            {'name': 'Brand-Sponsored Community', 'description': 'A space for customers to connect, share experiences, and engage directly with a brand.'},
        ],
        'Frontend Framework': [
            {'name': 'React.js', 'description': 'The library for web and native user interfaces.'},
            {'name': 'Next.js (React Framework)', 'description': 'A powerful React framework for building scalable, high-performance web applications.'},
            {'name': 'Nuxt.js (Vue Framework)', 'description': 'The intuitive Vue framework for creating universal applications and static websites.'},
            {'name': 'SvelteKit', 'description': 'A framework for building web applications with a beautiful development experience.'},
            {'name': 'Angular', 'description': 'A platform for building mobile and desktop web applications, developed by Google.'},
        ],
        'Backend & Database': [
            {'name': 'Supabase (Postgres)', 'description': 'An open-source Firebase alternative built on PostgreSQL with database, auth, and APIs.'},
            {'name': 'MySQL', 'description': 'The world\'s most popular open-source relational database.'},
            {'name': 'MongoDB (NoSQL)', 'description': 'A leading NoSQL, document-oriented database with flexible schema.'},
            {'name': 'Airtable', 'description': 'A low-code platform combining spreadsheet flexibility with database power.'},
            {'name': 'Self-Hosted Postgres + Node.js', 'description': 'Full control over backend infrastructure with Node.js application server.'},
            {'name': 'Laravel (PHP)', 'description': 'A PHP framework with expressive, elegant syntax for web applications.'},
            {'name': 'Firebase (NoSQL)', 'description': 'Google\'s platform with realtime NoSQL database, auth, and hosting.'},
        ],
        'Content Management (CMS)': [
            {'name': 'Strapi', 'description': 'Leading open-source headless CMS. 100% JavaScript, fully customizable.'},
            {'name': 'Payload CMS', 'description': 'TypeScript-native headless CMS that can install directly into Next.js apps.'},
            {'name': 'Headless WordPress', 'description': 'World\'s most popular CMS as a content backend via REST API.'},
            {'name': 'October CMS (Laravel)', 'description': 'Free, open-source CMS based on Laravel, known for simplicity and extensibility.'},
            {'name': 'Statamic (Laravel)', 'description': 'Flexible, flat-file CMS built on Laravel that stores content in files.'},
            {'name': 'Winter CMS (Laravel)', 'description': 'Community-driven fork of October CMS with focus on speed and security.'},
            {'name': 'Hygraph', 'description': 'Federated content platform to connect and manage content from various sources.'},
        ],
        'Primary AI Models': [
            {'name': 'Google Gemini', 'description': 'Natively multimodal models designed to understand text, images, audio, and video.'},
            {'name': 'Anthropic Claude', 'description': 'Renowned for strong performance in creative writing and complex reasoning.'},
            {'name': 'OpenAI GPT-4o', 'description': 'Highly advanced multimodal model capable of processing text, audio, and images.'},
            {'name': 'Ollama (Local LLM)', 'description': 'Run large language models locally on your own machine for privacy and control.'},
            {'name': 'LM Studio', 'description': 'Discover, download, and run local LLMs from a user-friendly desktop app.'},
        ],
        'AI Video Creation Tool': [
            {'name': 'HeyGen', 'description': 'AI video platform for creating engaging videos with generative avatars and voice cloning.'},
            {'name': 'Tavus', 'description': 'AI video personalization platform that generates unique videos for thousands of recipients.'},
            {'name': 'Pictory AI', 'description': 'Automatically creates short, shareable branded videos from long-form content.'},
            {'name': 'InVideo', 'description': 'Online video editor with templates, stock media, and AI tools.'},
            {'name': 'Google Veo', 'description': 'Google\'s most capable generative video model for high-quality, realistic video clips.'},
        ],
        'CRM & Marketing': [
            {'name': 'GoHighLevel', 'description': 'All-in-one sales and marketing platform for agencies and small businesses.'},
            {'name': 'HubSpot', 'description': 'Leading CRM platform with tools for marketing, sales, customer service, and content.'},
            {'name': 'Salesforce', 'description': 'The #1 AI CRM, uniting marketing, sales, service, commerce, and IT teams.'},
        ],
        'Payment Gateways': [
            {'name': 'Stripe', 'description': 'Developer-first payments platform with powerful APIs for transactions and subscriptions.'},
            {'name': 'PayPal', 'description': 'Globally recognized and trusted payment provider for online payments.'},
            {'name': 'BitPay', 'description': 'Leading crypto payment processor for accepting cryptocurrencies.'},
            {'name': 'Coinbase Commerce', 'description': 'Straightforward solution for businesses to accept cryptocurrency payments.'},
        ],
        'Hosting & Deployment': [
            {'name': 'Vercel (JS Platform)', 'description': 'Platform from Next.js creators with seamless git-based deployment.'},
            {'name': 'Netlify (JS Platform)', 'description': 'Powerful platform for building and deploying modern Jamstack web projects.'},
            {'name': 'Cloudflare Pages', 'description': 'Jamstack platform for frontend developers to collaborate and deploy websites.'},
            {'name': 'Render', 'description': 'Unified cloud to build and run apps with free TLS certificates and global CDN.'},
            {'name': 'Replit', 'description': 'Online IDE and hosting platform for rapid development in the browser.'},
            {'name': 'Hostinger (VPS)', 'description': 'Popular and affordable web hosting from shared hosting to VPS and cloud.'},
            {'name': 'Bluehost (VPS)', 'description': 'Leading web hosting solutions specializing in WordPress.'},
            {'name': 'SiteGround (VPS)', 'description': 'Known for top-rated customer support and fast, reliable hosting.'},
        ],
        'AI Site Creation': [
            {'name': 'Cursor', 'description': 'The AI-first code editor, designed for building software with LLMs at its core.'},
            {'name': 'VS Code', 'description': 'Free, powerful, and extensible code editor from Microsoft.'},
            {'name': 'Lovable', 'description': 'AI-powered platform to build web applications from plain English descriptions.'},
            {'name': 'Windsurf', 'description': 'AI tool for rapid site generation and development.'},
            {'name': 'Webstudio', 'description': 'Open-source visual development platform, alternative to Webflow.'},
            {'name': 'GrapesJS', 'description': 'Open-source Web Builder Framework for building HTML templates.'},
            {'name': 'DataButton', 'description': 'Build and deploy data apps, dashboards, and AI agents from your browser.'},
        ],
        'UI/UX Design Tools': [
            {'name': 'Figma', 'description': 'Leading collaborative interface design tool for creating, testing, and shipping designs.'},
            {'name': 'Relume', 'description': 'Build websites in hours with world\'s largest library of Figma & Webflow components.'},
            {'name': 'shadcn/ui', 'description': 'Beautifully designed components to copy and paste. Accessible and customizable.'},
            {'name': 'Lucide React', 'description': 'Beautiful and consistent open-source icon set as a React library.'},
            {'name': 'Anima', 'description': 'Design-to-code platform that exports Figma prototypes to HTML, React, or Vue.'},
        ],
        'Agentified Agent Army': [
            {'name': 'n8n', 'description': 'Powerful, self-hostable workflow automation tool to connect services.'},
            {'name': 'Zapier', 'description': 'Most popular no-code automation tool connecting thousands of web apps.'},
            {'name': 'Make', 'description': 'Visual platform to design, build, and automate tasks, workflows, and apps.'},
            {'name': 'LangGraph', 'description': 'Library for building stateful, multi-actor applications with LLMs.'},
            {'name': 'LangChain', 'description': 'Framework for developing applications powered by large language models.'},
            {'name': 'Lindy', 'description': 'AI-powered assistant for repetitive tasks like scheduling and email drafting.'},
            {'name': 'Relevance AI', 'description': 'Low-code platform for building and deploying AI agents and workflows.'},
        ],
        'MCP Servers': [
            {'name': 'Supabase MCP', 'description': 'Model-Context-Protocol server for interacting with Supabase databases.'},
            {'name': 'GitHub MCP', 'description': 'MCP server for AI agents to interact with GitHub repositories and issues.'},
            {'name': 'Docker MCP', 'description': 'MCP server for managing Docker containers through agentic interface.'},
            {'name': 'Crawl4AI MCP', 'description': 'MCP server for AI-powered web scraping and data retrieval for RAG pipelines.'},
            {'name': 'Brave Search MCP', 'description': 'Provides AI agents with access to Brave Search API for real-time web info.'},
            {'name': 'Figma MCP', 'description': 'MCP server for agentic interaction with Figma designs and prototypes.'},
            {'name': 'Google Drive MCP', 'description': 'MCP server for agents to read, write, and manage files in Google Drive.'},
            {'name': 'Playwright MCP', 'description': 'Allows agents to control Playwright browser for web automation.'},
        ],
        'Video Hosting Channel': [
            {'name': 'YouTube', 'description': 'World\'s largest video sharing platform for public content and brand channels.'},
            {'name': 'Vimeo', 'description': 'High-quality video platform favored by professionals with ad-free experience.'},
            {'name': 'Wistia (Marketing)', 'description': 'Focuses on video marketing with lead generation and analytics.'},
            {'name': 'Vidyard (Sales)', 'description': 'Designed for B2B sales with personalized video messages and tracking.'},
            {'name': 'Cloudinary (DAM)', 'description': 'Specializes in image and video scaling, optimization, and delivery.'},
            {'name': 'Orange Logic (DAM)', 'description': 'Best for handling large, high-resolution video files including 8K video.'},
            {'name': 'MediaValet (DAM)', 'description': 'Excels in automated video transcription and translation using AI.'},
            {'name': 'Bynder (DAM)', 'description': 'Strong in dynamic asset transformation and optimization.'},
        ],
        'Vector Database': [
            {'name': 'Pinecone', 'description': 'Managed, cloud-native vector database with simple API and no infrastructure hassles.'},
            {'name': 'Weaviate', 'description': 'Open-source, AI-native vector database for intuitive AI-powered applications.'},
            {'name': 'Milvus', 'description': 'Open-source vector database built for scalable similarity search for AI.'},
            {'name': 'Qdrant', 'description': 'Open-source vector similarity search engine written in Rust.'},
            {'name': 'ChromaDB', 'description': 'The open-source embedding database for building LLM apps with memory.'},
            {'name': 'PGVector', 'description': 'Open-source vector similarity search extension for PostgreSQL.'},
            {'name': 'Elasticsearch Vector Search', 'description': 'Leverages Elasticsearch for scalable vector search alongside text search.'},
        ],
    }

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide1, DARK_BG, prs)

    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "VidiSmart"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(72)
    title_p.font.bold = True
    title_p.font.color.rgb = ACCENT_GREEN
    title_p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI Site Planner"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(48)
    subtitle_p.font.color.rgb = WHITE
    subtitle_p.alignment = PP_ALIGN.CENTER

    tagline_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.6))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Configure your modern, AI SmartChannel Tech Stack"
    tagline_p = tagline_frame.paragraphs[0]
    tagline_p.font.size = Pt(24)
    tagline_p.font.color.rgb = GRAY
    tagline_p.alignment = PP_ALIGN.CENTER

    brands_box = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    brands_frame = brands_box.text_frame
    brands_frame.text = "VidiCRM™  •  VidiBlast & VidiMail™  •  VidiBot™"
    brands_p = brands_frame.paragraphs[0]
    brands_p.font.size = Pt(18)
    brands_p.font.color.rgb = YELLOW
    brands_p.alignment = PP_ALIGN.CENTER

    # Create a slide for each category
    for category_name, items in tech_data.items():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_dark_background(slide, DARK_BG, prs)
        add_slide_title(slide, category_name, ACCENT_GREEN)
        add_option_list(slide, items, WHITE)

    # Performance Metrics Slide
    perf_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(perf_slide, DARK_BG, prs)
    add_slide_title(perf_slide, "Performance Metrics", ACCENT_GREEN)

    perf_content = perf_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    perf_frame = perf_content.text_frame
    perf_frame.word_wrap = True

    perf_text = """VidiSmart.com delivers exceptional performance:

🚀 Largest Contentful Paint (LCP): 139ms
   Industry benchmark: < 2500ms = Good
   VidiSmart: 91% faster than "good" threshold

⚡ Time to First Byte (TTFB): 1ms
   Nearly instant server response

🎯 Cumulative Layout Shift (CLS): 0.00
   Perfect score - zero layout shifts

Technical Implementation:
• Canvas-based graphics (Three.js)
• No traditional image loading
• WebGL-accelerated animations
• Optimized JavaScript delivery"""

    perf_frame.text = perf_text
    for paragraph in perf_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(10)

    # Call to Action Slide
    cta_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(cta_slide, DARK_BG, prs)

    cta_title = cta_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    cta_title_frame = cta_title.text_frame
    cta_title_frame.text = "Ready to Build Your Stack?"
    cta_p = cta_title_frame.paragraphs[0]
    cta_p.font.size = Pt(54)
    cta_p.font.bold = True
    cta_p.font.color.rgb = ACCENT_GREEN
    cta_p.alignment = PP_ALIGN.CENTER

    visit_box = cta_slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    visit_frame = visit_box.text_frame
    visit_frame.text = "Visit: vidismart.com"
    visit_p = visit_frame.paragraphs[0]
    visit_p.font.size = Pt(36)
    visit_p.font.color.rgb = YELLOW
    visit_p.alignment = PP_ALIGN.CENTER

    footer_box = cta_slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.8))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Helping you to thrive in 2025\n©2025 VidiBuzz - https://VidiBuzz.com"
    for paragraph in footer_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = GRAY
        paragraph.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_path = '/mnt/m/code/vidismart/VidiSmart-AI-Site-Planner.pptx'
    prs.save(output_path)

    total_slides = 1 + len(tech_data) + 2  # Title + Categories + Performance + CTA
    print(f"✅ Presentation created successfully: {output_path}")
    print(f"📊 Total slides: {total_slides}")
    print(f"📋 Categories: {len(tech_data)}")
    print(f"🎨 Theme: VidiSmart Brand Colors (Dark mode with green accents)")
    print(f"\n📑 Slide breakdown:")
    print(f"   1. Title Slide")
    for i, category in enumerate(tech_data.keys(), 2):
        print(f"   {i}. {category}")
    print(f"   {total_slides-1}. Performance Metrics")
    print(f"   {total_slides}. Call to Action")

if __name__ == "__main__":
    create_vidismart_presentation()
